"""
pptx-exporter / export_pptx.py

v3 (Chromium-rendered layout)
-----------------------------
브라우저가 실제로 계산한 HTML/CSS 레이아웃을 Playwright/Chromium에서 읽고,
그 결과(x/y/width/height, computed style)를 python-pptx에 전달해 편집 가능한
네이티브 PPTX로 재구성한다.

핵심 변경:
- CSS Grid / Flexbox / % / transform을 직접 계산하지 않는다.
  Chromium의 getBoundingClientRect() 결과를 그대로 사용한다.
- getComputedStyle()로 최종 font/color/background/border/alignment를 읽는다.
- 텍스트는 직접 텍스트 노드 단위로 편집 가능한 PowerPoint TextBox로 만든다.
- img는 브라우저가 계산한 최종 위치/크기로 배치한다.
- svg/canvas는 Playwright가 해당 요소만 PNG로 캡처하여 정확한 위치에 배치한다.
- 배경색/테두리가 있는 DOM 요소는 PowerPoint 사각형으로 근사 재현한다.
- 기존 BeautifulSoup 기반 CSS 파서 / 세로 스택 fallback을 제거한다.

설치:
    pip install python-pptx pillow playwright
    playwright install chromium

사용:
    python export_pptx.py --web-ppt <web_ppt/vN 경로> --output <final.pptx 경로>

주의:
- HTML 렌더링 폰트와 PowerPoint 실행 PC의 설치 폰트가 같아야 줄바꿈 오차가 최소화된다.
- 복잡한 CSS gradient/filter/mask/clip-path는 네이티브 PowerPoint 도형으로 1:1 재현하기 어렵다.
  svg/canvas는 이미지로 보존한다.
"""

from __future__ import annotations

import argparse
import base64
import math
import re
import sys
import tempfile
from pathlib import Path
from urllib.parse import unquote, urlparse

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

SLIDE_WIDTH_PX = 1280
SLIDE_HEIGHT_PX = 720
EMU_PER_PX = 9525  # 96dpi
PX_TO_PT = 72.0 / 96.0
DEFAULT_FONT = "맑은 고딕"

# 너무 작은 장식/노이즈 박스는 무시
MIN_RENDER_W = 0.5
MIN_RENDER_H = 0.5


def px(v: float) -> Emu:
    return Emu(int(round(v * EMU_PER_PX)))


def px_to_pt(v: float) -> float:
    return float(v) * PX_TO_PT


class ConversionWarning(Exception):
    pass


# ---------------------------------------------------------------------
# CSS / color helpers
# ---------------------------------------------------------------------

_RGBA_RE = re.compile(
    r"rgba?\(\s*([\d.]+)\s*,\s*([\d.]+)\s*,\s*([\d.]+)"
    r"(?:\s*,\s*([\d.]+))?\s*\)",
    re.I,
)
_HEX_RE = re.compile(r"^#([0-9a-f]{3}|[0-9a-f]{6}|[0-9a-f]{8})$", re.I)


def parse_css_color(value: str | None) -> tuple[int, int, int, float] | None:
    """CSS color -> (r,g,b,alpha). transparent/invalid -> None."""
    if not value:
        return None
    s = value.strip().lower()
    if s in ("transparent", "rgba(0, 0, 0, 0)", "rgba(0,0,0,0)"):
        return None

    m = _RGBA_RE.match(s)
    if m:
        r = max(0, min(255, int(round(float(m.group(1))))))
        g = max(0, min(255, int(round(float(m.group(2))))))
        b = max(0, min(255, int(round(float(m.group(3))))))
        a = float(m.group(4)) if m.group(4) is not None else 1.0
        return r, g, b, max(0.0, min(1.0, a))

    m = _HEX_RE.match(s)
    if m:
        h = m.group(1)
        if len(h) == 3:
            r, g, b = (int(c * 2, 16) for c in h)
            return r, g, b, 1.0
        if len(h) == 6:
            return int(h[:2], 16), int(h[2:4], 16), int(h[4:6], 16), 1.0
        if len(h) == 8:
            return (
                int(h[:2], 16),
                int(h[2:4], 16),
                int(h[4:6], 16),
                int(h[6:8], 16) / 255.0,
            )
    return None


def set_fill(shape, css_color: str | None) -> bool:
    parsed = parse_css_color(css_color)
    if not parsed:
        shape.fill.background()
        return False
    r, g, b, a = parsed
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    # python-pptx transparency: 0 = opaque, 100000 = fully transparent
    try:
        shape.fill.transparency = int(round((1.0 - a) * 100000))
    except Exception:
        pass
    return True


def set_line(shape, css_color: str | None, width_px: float = 0.0) -> None:
    parsed = parse_css_color(css_color)
    if not parsed or width_px <= 0:
        try:
            shape.line.fill.background()
        except Exception:
            pass
        return
    r, g, b, a = parsed
    shape.line.color.rgb = RGBColor(r, g, b)
    shape.line.width = Pt(max(0.1, px_to_pt(width_px)))
    try:
        shape.line.transparency = int(round((1.0 - a) * 100000))
    except Exception:
        pass


_GRADIENT_RE = re.compile(r"linear-gradient\(\s*(.*)\)\s*$", re.I | re.S)
_ANGLE_KEYWORDS = {
    "to top": 0.0, "to right": 90.0, "to bottom": 180.0, "to left": 270.0,
    "to top right": 45.0, "to right top": 45.0,
    "to bottom right": 135.0, "to right bottom": 135.0,
    "to bottom left": 225.0, "to left bottom": 225.0,
    "to top left": 315.0, "to left top": 315.0,
}


def _split_top_level(s: str) -> list[str]:
    """괄호(rgba(...) 등) 안의 콤마는 무시하고 최상위 콤마로만 분리한다."""
    parts, depth, buf = [], 0, []
    for ch in s:
        if ch == "(":
            depth += 1
        elif ch == ")":
            depth -= 1
        if ch == "," and depth == 0:
            parts.append("".join(buf))
            buf = []
        else:
            buf.append(ch)
    if buf:
        parts.append("".join(buf))
    return [p.strip() for p in parts if p.strip()]


def parse_linear_gradient(value: str | None) -> tuple[float, list[tuple[float, int, int, int, float]]] | None:
    """CSS `linear-gradient(...)` 문자열을 (각도(deg, CSS 기준), [(위치 0..1, r,g,b,alpha), ...])로 변환한다.
    변환 대상이 아니면 None을 반환한다(다른 gradient 종류, 파싱 실패 등)."""
    if not value:
        return None
    m = _GRADIENT_RE.search(value)
    if not m:
        return None
    parts = _split_top_level(m.group(1))
    if not parts:
        return None

    angle_deg = 180.0  # CSS 기본값: to bottom
    first = parts[0].strip().lower()
    stop_parts = parts
    deg_m = re.match(r"^(-?[\d.]+)deg$", first)
    if deg_m:
        angle_deg = float(deg_m.group(1)) % 360.0
        stop_parts = parts[1:]
    elif first in _ANGLE_KEYWORDS:
        angle_deg = _ANGLE_KEYWORDS[first]
        stop_parts = parts[1:]
    elif first.startswith("to "):
        angle_deg = _ANGLE_KEYWORDS.get(first, 180.0)
        stop_parts = parts[1:]

    stops: list[tuple[float, int, int, int, float]] = []
    n = len(stop_parts)
    for i, sp in enumerate(stop_parts):
        sp = sp.strip()
        pos_m = re.search(r"(-?[\d.]+)%\s*$", sp)
        color_str = sp[: pos_m.start()].strip() if pos_m else sp
        parsed = parse_css_color(color_str)
        if not parsed:
            continue
        r, g, b, a = parsed
        if pos_m:
            pos = max(0.0, min(1.0, float(pos_m.group(1)) / 100.0))
        else:
            pos = 0.0 if n <= 1 else i / (n - 1)
        stops.append((pos, r, g, b, a))

    if len(stops) < 2:
        return None
    stops.sort(key=lambda s: s[0])
    return angle_deg, stops


def apply_linear_gradient(shape, angle_css_deg: float,
                           stops: list[tuple[float, int, int, int, float]]) -> None:
    """python-pptx 공개 API는 그라데이션 stop을 기본 2개로 제한하므로,
    <a:gradFill>의 <a:gsLst>를 직접 조작해 임의 개수의 stop을 반영한다."""
    shape.fill.gradient()
    spPr = shape.fill._xPr
    grad_fill = spPr.find(qn("a:gradFill"))
    if grad_fill is None:
        return
    gs_lst = grad_fill.find(qn("a:gsLst"))
    if gs_lst is None:
        return
    for child in list(gs_lst):
        gs_lst.remove(child)

    for pos, r, g, b, a in stops:
        gs = etree.SubElement(gs_lst, qn("a:gs"))
        gs.set("pos", str(int(round(pos * 100000))))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", f"{r:02X}{g:02X}{b:02X}")
        if a < 1.0:
            alpha_el = etree.SubElement(srgb, qn("a:alpha"))
            alpha_el.set("val", str(int(round(a * 100000))))

    lin = grad_fill.find(qn("a:lin"))
    if lin is None:
        lin = etree.SubElement(grad_fill, qn("a:lin"))
    ppt_angle = (angle_css_deg - 90.0) % 360.0  # CSS 0deg(위쪽) -> PPT 0deg(오른쪽) 기준 보정
    lin.set("ang", str(int(round(ppt_angle * 60000)) % 21600000))
    lin.set("scaled", "1")


def css_font_name(font_family: str | None) -> str:
    if not font_family:
        return DEFAULT_FONT
    # 첫 번째 family를 우선
    first = font_family.split(",")[0].strip().strip("'\"")
    return first or DEFAULT_FONT


def css_weight_is_bold(weight: str | None) -> bool:
    if not weight:
        return False
    if weight.lower() in ("bold", "bolder"):
        return True
    try:
        return int(float(weight)) >= 600
    except Exception:
        return False


def parse_px_string(value: str | None, default: float = 0.0) -> float:
    if not value:
        return default
    m = re.match(r"^\s*(-?[\d.]+)px\s*$", str(value))
    return float(m.group(1)) if m else default


def is_visible_item(item: dict) -> bool:
    return (
        item.get("visible", True)
        and float(item.get("w", 0)) >= MIN_RENDER_W
        and float(item.get("h", 0)) >= MIN_RENDER_H
    )


# ---------------------------------------------------------------------
# Browser rendering / extraction
# ---------------------------------------------------------------------

EXTRACT_JS = r"""
async () => {
  await document.fonts.ready;

  const slides = Array.from(document.querySelectorAll('.slide'));

  const rgbaAlpha = (s) => {
    if (!s || s === 'transparent') return 0;
    const m = s.match(/rgba?\(([^)]+)\)/i);
    if (!m) return 1;
    const parts = m[1].split(',').map(x => x.trim());
    return parts.length >= 4 ? Number(parts[3]) : 1;
  };

  const num = (s) => {
    const v = parseFloat(s || '0');
    return Number.isFinite(v) ? v : 0;
  };

  const isVisible = (el, cs, r) => {
    if (!r || r.width <= 0 || r.height <= 0) return false;
    if (cs.display === 'none' || cs.visibility === 'hidden') return false;
    if (Number(cs.opacity || 1) <= 0) return false;
    return true;
  };

  const directTextNodes = (el) =>
    Array.from(el.childNodes).filter(n => n.nodeType === Node.TEXT_NODE && n.nodeValue.trim());

  const textInfo = (node, slideRect, order, parentEl) => {
    const range = document.createRange();
    range.selectNodeContents(node);
    const rr = range.getBoundingClientRect();
    const cs = getComputedStyle(parentEl);

    let text = node.nodeValue;
    // HTML normal whitespace behavior를 PPT용으로 안정화.
    if (cs.whiteSpace === 'normal' || cs.whiteSpace === 'nowrap') {
      text = text.replace(/\s+/g, ' ').trim();
    } else {
      text = text.replace(/\r/g, '').trim();
    }

    return {
      kind: 'text',
      order,
      text,
      x: rr.left - slideRect.left,
      y: rr.top - slideRect.top,
      w: rr.width,
      h: rr.height,
      visible: rr.width > 0 && rr.height > 0 && !!text,
      style: {
        fontFamily: cs.fontFamily,
        fontSize: cs.fontSize,
        fontWeight: cs.fontWeight,
        fontStyle: cs.fontStyle,
        color: cs.color,
        textAlign: cs.textAlign,
        lineHeight: cs.lineHeight,
        letterSpacing: cs.letterSpacing,
        textDecorationLine: cs.textDecorationLine,
        whiteSpace: cs.whiteSpace,
        direction: cs.direction
      }
    };
  };

  const all = [];

  for (let si = 0; si < slides.length; si++) {
    const slide = slides[si];
    const sr = slide.getBoundingClientRect();

    // .slide가 1280x720이 아닌 경우도 실제 렌더링 결과를 기준으로 scale 가능하도록 기록
    const result = {
      slideIndex: si,
      pageX: sr.left + window.scrollX,
      pageY: sr.top + window.scrollY,
      renderedW: sr.width,
      renderedH: sr.height,
      items: []
    };

    let order = 0;

    const walker = document.createTreeWalker(slide, NodeFilter.SHOW_ELEMENT);
    let el = slide;

    while (el) {
      const cs = getComputedStyle(el);
      const r = el.getBoundingClientRect();
      const localX = r.left - sr.left;
      const localY = r.top - sr.top;
      const visible = isVisible(el, cs, r);

      // complex element screenshot identifier
      if ((el.tagName === 'SVG' || el.tagName === 'CANVAS') && visible) {
        const exportId = `ppt-export-${si}-${order}-${Math.random().toString(36).slice(2)}`;
        el.setAttribute('data-ppt-export-id', exportId);
        result.items.push({
          kind: 'complex',
          order: order++,
          exportId,
          tag: el.tagName.toLowerCase(),
          x: localX, y: localY, w: r.width, h: r.height,
          visible
        });

        // SVG/canvas 내부 텍스트를 중복 native text로 올리지 않도록 하위 탐색 생략은
        // TreeWalker에서 직접 제어하기 까다로워 later filtering에서 complex ancestor를 체크한다.
      } else if (el.tagName === 'TABLE' && visible) {
        // <table>은 편집 가능한 네이티브 PowerPoint 표로 재구성하기 위해
        // 행/셀 구조를 통째로 캡처한다(자손은 아래에서 별도 box/text로 중복 캡처하지 않음).
        const csTable = getComputedStyle(el);
        const rowEls = Array.from(el.querySelectorAll('tr'));
        const rows = rowEls.map((tr, rowIdx) => {
          const csTr = getComputedStyle(tr);
          return Array.from(tr.children).filter(c => c.tagName === 'TD' || c.tagName === 'TH').map(td => {
            const cr = td.getBoundingClientRect();
            const csTd = getComputedStyle(td);
            return {
              text: td.textContent.replace(/\s+/g, ' ').trim(),
              isHeader: td.tagName === 'TH',
              x: cr.left - sr.left, y: cr.top - sr.top, w: cr.width, h: cr.height,
              style: {
                backgroundColor: csTd.backgroundColor,
                color: csTd.color,
                fontFamily: csTd.fontFamily,
                fontSize: csTd.fontSize,
                fontWeight: csTd.fontWeight,
                textAlign: csTd.textAlign,
                // 표 자체의 top border(첫 행에만), 각 행의 bottom border
                borderTopWidth: rowIdx === 0 ? num(csTable.borderTopWidth) : 0,
                borderTopColor: csTable.borderTopColor,
                borderBottomWidth: num(csTr.borderBottomWidth),
                borderBottomColor: csTr.borderBottomColor
              }
            };
          });
        });
        result.items.push({
          kind: 'table',
          order: order++,
          x: localX, y: localY, w: r.width, h: r.height,
          visible,
          rows
        });
      } else {
        const bg = cs.backgroundColor;
        const borderWidths = {
          top: num(cs.borderTopWidth),
          right: num(cs.borderRightWidth),
          bottom: num(cs.borderBottomWidth),
          left: num(cs.borderLeftWidth)
        };
        const hasBorder = Object.values(borderWidths).some(v => v > 0);
        const hasBg = rgbaAlpha(bg) > 0;
        const bgImage = cs.backgroundImage && cs.backgroundImage !== 'none' ? cs.backgroundImage : null;
        const insideTable = !!el.closest('table');

        if (visible && !insideTable && (hasBg || hasBorder || bgImage)) {
          result.items.push({
            kind: 'box',
            order: order++,
            tag: el.tagName.toLowerCase(),
            x: localX, y: localY, w: r.width, h: r.height,
            visible,
            style: {
              backgroundColor: bg,
              backgroundImage: bgImage,
              backgroundSize: cs.backgroundSize,
              backgroundPosition: cs.backgroundPosition,
              backgroundRepeat: cs.backgroundRepeat,
              borderTopWidth: borderWidths.top,
              borderRightWidth: borderWidths.right,
              borderBottomWidth: borderWidths.bottom,
              borderLeftWidth: borderWidths.left,
              borderTopColor: cs.borderTopColor,
              borderRightColor: cs.borderRightColor,
              borderBottomColor: cs.borderBottomColor,
              borderLeftColor: cs.borderLeftColor,
              borderRadius: cs.borderRadius,
              opacity: cs.opacity
            }
          });
        }

        if (el.tagName === 'IMG' && visible) {
          result.items.push({
            kind: 'image',
            order: order++,
            x: localX, y: localY, w: r.width, h: r.height,
            visible,
            src: el.currentSrc || el.src || '',
            style: {
              objectFit: cs.objectFit,
              objectPosition: cs.objectPosition,
              opacity: cs.opacity
            }
          });
        }

        // SVG/canvas의 자손 텍스트는 캡처 이미지에 이미 포함되고, table의 자손 텍스트는
        // 위 table 캡처에서 셀 단위로 이미 포함되므로 중복 native text를 만들지 않는다.
        const hasComplexAncestor = !!el.closest('svg,canvas,table');
        if (!hasComplexAncestor) {
          for (const tn of directTextNodes(el)) {
            const ti = textInfo(tn, sr, order++, el);
            if (ti.visible) result.items.push(ti);
          }
        }
      }

      el = walker.nextNode();
    }

    all.push(result);
  }

  return all;
}
"""


def ensure_playwright():
    try:
        from playwright.sync_api import sync_playwright
        return sync_playwright
    except Exception as e:
        raise SystemExit(
            "Playwright가 필요합니다.\n"
            "  pip install playwright\n"
            "  playwright install chromium\n"
            f"원인: {e}"
        )


def render_html(web_ppt_dir: Path, temp_dir: Path) -> tuple[list[dict], dict[str, Path]]:
    index_path = (web_ppt_dir / "index.html").resolve()
    if not index_path.exists():
        raise SystemExit(f"index.html을 찾을 수 없습니다: {index_path}")

    sync_playwright = ensure_playwright()
    complex_images: dict[str, Path] = {}

    with sync_playwright() as p:
        try:
            browser = p.chromium.launch()
        except Exception as e:
            raise SystemExit(
                "Chromium을 실행하지 못했습니다. 다음 명령을 먼저 실행하세요:\n"
                "  playwright install chromium\n"
                f"원인: {e}"
            )

        page = browser.new_page(
            viewport={"width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX},
            device_scale_factor=1,
        )

        # 로컬 file:// 리소스 사용
        page.goto(index_path.as_uri(), wait_until="load")
        page.evaluate("() => document.fonts.ready")
        # JS/Chart 등이 한 프레임 이상 렌더링될 시간을 확보
        page.wait_for_timeout(250)

        data = page.evaluate(EXTRACT_JS)

        # SVG / canvas는 실제 렌더링 결과를 PNG로 캡처
        for slide in data:
            for item in slide["items"]:
                if item.get("kind") != "complex":
                    continue
                export_id = item["exportId"]
                out = temp_dir / f"{export_id}.png"
                loc = page.locator(f'[data-ppt-export-id="{export_id}"]')
                try:
                    loc.screenshot(path=str(out), omit_background=True)
                    complex_images[export_id] = out
                except Exception:
                    # omit_background 미지원/실패 시 일반 screenshot 재시도
                    try:
                        loc.screenshot(path=str(out))
                        complex_images[export_id] = out
                    except Exception:
                        pass

        browser.close()

    return data, complex_images


# ---------------------------------------------------------------------
# Asset / image helpers
# ---------------------------------------------------------------------

_URL_RE = re.compile(r'url\(["\']?([^"\')]+)["\']?\)', re.I)


def css_url(background_image: str | None) -> str | None:
    if not background_image:
        return None
    m = _URL_RE.search(background_image)
    return m.group(1) if m else None


def resolve_src_to_path(src: str, web_ppt_dir: Path, temp_dir: Path) -> Path | None:
    if not src:
        return None

    if src.startswith("data:image/"):
        m = re.match(r"data:image/([^;]+);base64,(.+)", src, re.S)
        if not m:
            return None
        ext = m.group(1).lower().replace("jpeg", "jpg").replace("svg+xml", "svg")
        out = temp_dir / f"data-image-{abs(hash(src))}.{ext}"
        try:
            out.write_bytes(base64.b64decode(m.group(2)))
            return out
        except Exception:
            return None

    parsed = urlparse(src)
    if parsed.scheme == "file":
        # url2pathname은 'file:///C:/...' 같은 Windows 드라이브 경로의 선행 슬래시를
        # 올바르게 제거해준다(단순 unquote만 하면 '/C:/...'가 되어 존재하지 않는 경로가 됨).
        from urllib.request import url2pathname
        return Path(url2pathname(parsed.path)).resolve()

    if parsed.scheme in ("http", "https"):
        # 브라우저는 표시할 수 있어도 python-pptx가 URL을 직접 삽입할 수는 없다.
        # 외부 URL 이미지는 현재 네이티브 삽입 대상에서 제외.
        return None

    p = Path(unquote(src))
    if p.is_absolute():
        return p
    return (web_ppt_dir / p).resolve()


def add_picture_exact(slide, img_path: Path, x: float, y: float, w: float, h: float, warnings: list[str]) -> None:
    if not img_path.exists():
        warnings.append(f"이미지 파일을 찾을 수 없어 건너뜀: {img_path}")
        return
    try:
        slide.shapes.add_picture(str(img_path), px(x), px(y), width=px(w), height=px(h))
    except Exception as e:
        warnings.append(f"이미지 삽입 실패: {img_path} ({e})")


# ---------------------------------------------------------------------
# PPT shapes
# ---------------------------------------------------------------------

def _thin_line_rect(slide, x: float, y: float, w: float, h: float, css_color: str | None) -> None:
    parsed = parse_css_color(css_color)
    if not parsed or w <= 0 or h <= 0:
        return
    r, g, b, a = parsed
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(max(w, 0.5)), px(max(h, 0.5)))
    line.shadow.inherit = False
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(r, g, b)
    try:
        line.fill.transparency = int(round((1.0 - a) * 100000))
    except Exception:
        pass


def add_side_borders(slide, x: float, y: float, w: float, h: float, style: dict, sx: float, sy: float) -> None:
    """CSS의 top/right/bottom/left border를 각각 독립된 얇은 선(사각형)으로 그린다
    (하나의 outline으로 근사하면 없는 변까지 테두리가 생기므로 변마다 따로 그린다)."""
    tw = float(style.get("borderTopWidth", 0) or 0) * sy
    rw = float(style.get("borderRightWidth", 0) or 0) * sx
    bw = float(style.get("borderBottomWidth", 0) or 0) * sy
    lw = float(style.get("borderLeftWidth", 0) or 0) * sx

    if tw > 0:
        _thin_line_rect(slide, x, y, w, tw, style.get("borderTopColor"))
    if bw > 0:
        _thin_line_rect(slide, x, y + h - bw, w, bw, style.get("borderBottomColor"))
    if lw > 0:
        _thin_line_rect(slide, x, y, lw, h, style.get("borderLeftColor"))
    if rw > 0:
        _thin_line_rect(slide, x + w - rw, y, rw, h, style.get("borderRightColor"))


def add_box(slide, item: dict, sx: float, sy: float, web_ppt_dir: Path, temp_dir: Path,
            warnings: list[str]) -> None:
    x = item["x"] * sx
    y = item["y"] * sy
    w = item["w"] * sx
    h = item["h"] * sy
    style = item.get("style", {})

    if w <= 0 or h <= 0:
        return

    # CSS background-image가 단순 url()이면 이미지로 우선 처리
    bg_src = css_url(style.get("backgroundImage"))
    if bg_src:
        p = resolve_src_to_path(bg_src, web_ppt_dir, temp_dir)
        if p and p.exists():
            # background-size/position의 cover/contain을 완벽히 native crop 재현하기보다는
            # 일단 렌더링 bbox에 맞춰 삽입. 배경색 overlay는 아래 사각형으로 같이 유지.
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))
            set_fill(shape, style.get("backgroundColor"))
            set_line(shape, None, 0)
            add_picture_exact(slide, p, x, y, w, h, warnings)
            return

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(x), px(y), px(w), px(h))

    gradient = parse_linear_gradient(style.get("backgroundImage"))
    if gradient:
        angle, stops = gradient
        apply_linear_gradient(shape, angle, stops)
    else:
        set_fill(shape, style.get("backgroundColor"))

    # 이 사각형 자체는 outline을 갖지 않는다. CSS는 border-left 하나만 있는 경우가 흔한데
    # 사각형 전체에 하나의 outline을 씌우면 존재하지 않는 3면까지 테두리가 생겨버린다.
    # 실제로 두께가 있는 변만, 그 변의 위치에 얇은 선(사각형)으로 개별 렌더링한다.
    shape.line.fill.background()
    add_side_borders(slide, x, y, w, h, style, sx, sy)


def add_text(slide, item: dict, sx: float, sy: float, warnings: list[str]) -> None:
    text = item.get("text", "")
    if not text:
        return

    x = item["x"] * sx
    y = item["y"] * sy
    w = max(item["w"] * sx, 2.0)
    h = max(item["h"] * sy, 2.0)
    style = item.get("style", {})

    box = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = style.get("whiteSpace") != "nowrap"
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = text
    p.space_before = Pt(0)
    p.space_after = Pt(0)

    align = (style.get("textAlign") or "").lower()
    if align in ("center", "-webkit-center"):
        p.alignment = PP_ALIGN.CENTER
    elif align in ("right", "end"):
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    line_height_px = parse_px_string(style.get("lineHeight"), 0.0)
    if line_height_px > 0:
        try:
            p.line_spacing = Pt(px_to_pt(line_height_px * sy))
        except Exception:
            pass

    # text node 하나이므로 보통 run 하나
    for run in p.runs:
        font = run.font
        fs_px = parse_px_string(style.get("fontSize"), 16.0)
        font.size = Pt(max(1.0, px_to_pt(fs_px * sy)))
        font.bold = css_weight_is_bold(style.get("fontWeight"))
        font.italic = (style.get("fontStyle") or "").lower() == "italic"
        font.name = css_font_name(style.get("fontFamily"))

        color = parse_css_color(style.get("color"))
        if color:
            r, g, b, a = color
            font.color.rgb = RGBColor(r, g, b)
            # font transparency는 python-pptx public API에서 직접 지원하지 않음

        deco = (style.get("textDecorationLine") or "").lower()
        if "underline" in deco:
            try:
                font.underline = True
            except Exception:
                pass


_TCPR_CHILD_ORDER = [
    "a:lnL", "a:lnR", "a:lnT", "a:lnB", "a:lnTlToBr", "a:lnBlToTr", "a:cell3D",
    "a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill", "a:grpFill",
]


def set_cell_border(cell, side: str, width_emu, css_color: str | None) -> None:
    """side: 'T'/'B'/'L'/'R'. tcPr 안에서 lnX는 반드시 fill류보다 앞에 와야 하므로
    스키마 순서를 지켜 올바른 위치에 삽입한다."""
    parsed = parse_css_color(css_color)
    if not parsed or width_emu <= 0:
        return
    r, g, b, a = parsed
    tag = f"a:ln{side}"
    tcPr = cell._tc.get_or_add_tcPr()
    existing = tcPr.find(qn(tag))
    if existing is not None:
        tcPr.remove(existing)
    ln = etree.Element(qn(tag))
    ln.set("w", str(int(width_emu)))
    ln.set("cap", "flat")
    solid_fill = etree.SubElement(ln, qn("a:solidFill"))
    srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
    srgb.set("val", f"{r:02X}{g:02X}{b:02X}")

    my_idx = _TCPR_CHILD_ORDER.index(tag)
    insert_at = len(list(tcPr))
    for i, child in enumerate(tcPr):
        child_tag = "a:" + etree.QName(child).localname
        if child_tag in _TCPR_CHILD_ORDER and _TCPR_CHILD_ORDER.index(child_tag) > my_idx:
            insert_at = i
            break
    tcPr.insert(insert_at, ln)


def add_table_exact(slide, item: dict, sx: float, sy: float, warnings: list[str]) -> None:
    """브라우저가 계산한 실제 셀 좌표를 바탕으로, 행/열 구조가 살아있는
    편집 가능한 네이티브 PowerPoint 표를 만든다(사각형+텍스트 조합이 아님)."""
    rows = item.get("rows") or []
    rows = [r for r in rows if r]
    if not rows:
        return

    n_rows = len(rows)
    n_cols = max(len(r) for r in rows)
    if n_cols == 0:
        return

    x = item["x"] * sx
    y = item["y"] * sy
    w = max(item["w"] * sx, 10.0)
    h = max(item["h"] * sy, 10.0)

    gframe = slide.shapes.add_table(n_rows, n_cols, px(x), px(y), px(w), px(h))
    table = gframe.table

    # 열 폭: 첫 행 각 셀의 실측 폭을 사용(비율 유지), 마지막 열은 나머지로 보정해 합을 정확히 맞춘다.
    first_row = rows[0]
    col_widths = [max(c["w"] * sx, 4.0) for c in first_row]
    while len(col_widths) < n_cols:
        col_widths.append(w / n_cols)
    total_w = sum(col_widths)
    col_widths = [cw * (w / total_w) for cw in col_widths] if total_w > 0 else col_widths
    for c_idx, cw in enumerate(col_widths):
        table.columns[c_idx].width = px(cw)

    # 행 높이: 각 행의 실측 높이를 사용, 마지막 행은 나머지로 보정.
    row_heights = [max(r[0]["h"] * sy, 4.0) if r else h / n_rows for r in rows]
    total_h = sum(row_heights)
    row_heights = [rh * (h / total_h) for rh in row_heights] if total_h > 0 else row_heights
    for r_idx, rh in enumerate(row_heights):
        table.rows[r_idx].height = px(rh)

    for r_idx, row in enumerate(rows):
        for c_idx in range(n_cols):
            if c_idx >= len(row):
                continue
            cell_data = row[c_idx]
            cell = table.cell(r_idx, c_idx)
            cell.margin_left = px(4)
            cell.margin_right = px(4)
            cell.margin_top = px(2)
            cell.margin_bottom = px(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            style = cell_data.get("style", {})

            top_w = float(style.get("borderTopWidth", 0) or 0) * sy
            if top_w > 0:
                set_cell_border(cell, "T", px(top_w), style.get("borderTopColor"))
            bottom_w = float(style.get("borderBottomWidth", 0) or 0) * sy
            if bottom_w > 0:
                set_cell_border(cell, "B", px(bottom_w), style.get("borderBottomColor"))

            bgcolor = parse_css_color(style.get("backgroundColor"))
            if bgcolor:
                r, g, b, a = bgcolor
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(r, g, b)
            else:
                cell.fill.background()

            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell_data.get("text", "")

            align = (style.get("textAlign") or "").lower()
            if align in ("center", "-webkit-center"):
                p.alignment = PP_ALIGN.CENTER
            elif align in ("right", "end"):
                p.alignment = PP_ALIGN.RIGHT
            else:
                p.alignment = PP_ALIGN.LEFT

            for run in p.runs:
                font = run.font
                fs_px = parse_px_string(style.get("fontSize"), 14.0)
                font.size = Pt(max(1.0, px_to_pt(fs_px * sy)))
                font.bold = css_weight_is_bold(style.get("fontWeight")) or bool(cell_data.get("isHeader"))
                font.name = css_font_name(style.get("fontFamily"))
                color = parse_css_color(style.get("color"))
                if color:
                    cr, cg, cb, _ = color
                    font.color.rgb = RGBColor(cr, cg, cb)

    # 표 자체 스타일(줄무늬/헤더 강조 등 기본 테마)이 원본과 다르게 보이는 것을 막기 위해
    # 첫 행 강조와 줄무늬를 모두 끈다 — 셀 배경색은 위에서 이미 원본 값대로 지정했다.
    tbl_el = table._tbl
    tblPr = tbl_el.find(qn("a:tblPr"))
    if tblPr is not None:
        tblPr.set("firstRow", "0")
        tblPr.set("bandRow", "0")


def build_slide(prs: Presentation, slide_data: dict, web_ppt_dir: Path, temp_dir: Path,
                complex_images: dict[str, Path], warnings: list[str]) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    rendered_w = float(slide_data.get("renderedW") or SLIDE_WIDTH_PX)
    rendered_h = float(slide_data.get("renderedH") or SLIDE_HEIGHT_PX)
    sx = SLIDE_WIDTH_PX / rendered_w if rendered_w else 1.0
    sy = SLIDE_HEIGHT_PX / rendered_h if rendered_h else 1.0

    items = sorted(slide_data.get("items", []), key=lambda x: x.get("order", 0))

    # DOM 순서대로 shape를 쌓는다. 부모 background -> 자식 -> text 순서가 대체로 유지된다.
    for item in items:
        if not is_visible_item(item):
            continue

        kind = item.get("kind")
        try:
            if kind == "box":
                add_box(slide, item, sx, sy, web_ppt_dir, temp_dir, warnings)

            elif kind == "image":
                src = item.get("src", "")
                p = resolve_src_to_path(src, web_ppt_dir, temp_dir)
                if p and p.exists():
                    add_picture_exact(
                        slide,
                        p,
                        item["x"] * sx,
                        item["y"] * sy,
                        item["w"] * sx,
                        item["h"] * sy,
                        warnings,
                    )
                else:
                    warnings.append(f"이미지 소스를 로컬 파일로 해석하지 못해 건너뜀: {src}")

            elif kind == "complex":
                p = complex_images.get(item.get("exportId", ""))
                if p and p.exists():
                    add_picture_exact(
                        slide,
                        p,
                        item["x"] * sx,
                        item["y"] * sy,
                        item["w"] * sx,
                        item["h"] * sy,
                        warnings,
                    )
                else:
                    warnings.append(
                        f"{item.get('tag', 'complex')} 캡처 실패: {item.get('exportId', '')}"
                    )

            elif kind == "text":
                add_text(slide, item, sx, sy, warnings)

            elif kind == "table":
                add_table_exact(slide, item, sx, sy, warnings)

        except Exception as e:
            warnings.append(f"요소 변환 실패(kind={kind}): {e}")


def convert(web_ppt_dir: Path, output_path: Path) -> list[str]:
    warnings: list[str] = []
    web_ppt_dir = web_ppt_dir.resolve()

    with tempfile.TemporaryDirectory(prefix="pptx-exporter-") as td:
        temp_dir = Path(td)
        slide_data, complex_images = render_html(web_ppt_dir, temp_dir)

        if not slide_data:
            raise SystemExit("index.html에서 .slide 요소를 찾지 못했습니다.")

        prs = Presentation()
        prs.slide_width = px(SLIDE_WIDTH_PX)
        prs.slide_height = px(SLIDE_HEIGHT_PX)

        for sd in slide_data:
            build_slide(prs, sd, web_ppt_dir, temp_dir, complex_images, warnings)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

    # 기본 검증
    reopened = Presentation(str(output_path))
    if len(reopened.slides) != len(slide_data):
        raise ConversionWarning(
            f"슬라이드 수 불일치: 원본 {len(slide_data)}개, 변환 결과 {len(reopened.slides)}개"
        )

    return warnings


def main() -> None:
    parser = argparse.ArgumentParser(
        description="Chromium 렌더링 결과를 기반으로 웹PPT를 편집 가능한 .pptx로 변환한다."
    )
    parser.add_argument("--web-ppt", required=True, help="web_ppt/vN 폴더 경로")
    parser.add_argument("--output", required=True, help="출력 .pptx 경로")
    args = parser.parse_args()

    web_ppt_dir = Path(args.web_ppt)
    output_path = Path(args.output)

    try:
        warnings = convert(web_ppt_dir, output_path)
    except ConversionWarning as e:
        print(f"[pptx-exporter] 검증 실패: {e}", file=sys.stderr)
        sys.exit(1)

    print(f"[pptx-exporter] 변환 완료 -> {output_path}")
    if warnings:
        print(f"[pptx-exporter] 경고 {len(warnings)}건:")
        for w in warnings:
            print(f"  - {w}")


if __name__ == "__main__":
    main()
