"""
web-ppt-generator / analyze_reference.py

레퍼런스 PPT(.pptx)의 구조(레이아웃, 색상, 폰트, 요소 배치)를 덤프해
LLM이 가변 규칙 후보를 해석할 수 있는 원재료를 만든다.

주의: 이 스크립트는 데이터를 추출만 한다. "좋은 디자인 요소"에 대한 해석·
가변 규칙 제안은 LLM(content-designer)의 판단 영역이며, 이 스크립트는 그
판단의 입력값만 제공한다. design-rules.md에 대한 반영 여부는 항상 사용자
승인을 거친다.

사용법:
    python analyze_reference.py --input <레퍼런스.pptx> [--output <json 경로>]
    출력 생략 시 표준출력으로 JSON을 인쇄한다.
"""
from __future__ import annotations

import argparse
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

EMU_PER_INCH = 914400


def emu_to_in(v) -> float:
    if v is None:
        return None
    return round(int(v) / EMU_PER_INCH, 3)


def color_to_hex(color) -> str | None:
    try:
        if color.type is None:
            return None
        if color.type == 1:  # MSO_THEME_COLOR / RGB depends on type value; guard broadly
            pass
        return str(color.rgb)
    except Exception:
        try:
            return f"theme:{color.theme_color}"
        except Exception:
            return None


def extract_theme_colors(prs: Presentation) -> dict:
    """slide master의 theme1.xml에서 색상 스킴(a:clrScheme)을 추출한다."""
    colors = {}
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clr_scheme = theme_part._element.find(".//a:clrScheme", ns)
        if clr_scheme is not None:
            for child in clr_scheme:
                tag = child.tag.split("}")[-1]
                srgb = child.find("a:srgbClr", ns)
                sys_clr = child.find("a:sysClr", ns)
                if srgb is not None:
                    colors[tag] = "#" + srgb.get("val")
                elif sys_clr is not None:
                    colors[tag] = "#" + sys_clr.get("lastClr", sys_clr.get("val", ""))
    except Exception as e:
        colors["_error"] = f"테마 색상 추출 실패: {e}"
    return colors


def extract_font_info(text_frame) -> list[dict]:
    fonts = []
    for para in text_frame.paragraphs:
        for run in para.runs:
            if not run.text.strip():
                continue
            f = run.font
            fonts.append({
                "text_sample": run.text.strip()[:40],
                "name": f.name,
                "size_pt": f.size.pt if f.size else None,
                "bold": f.bold,
                "italic": f.italic,
                "color": color_to_hex(f.color) if f.color else None,
                "alignment": str(para.alignment) if para.alignment else None,
            })
    return fonts


def extract_shape(shape) -> dict:
    info = {
        "shape_type": str(shape.shape_type),
        "name": shape.name,
        "left_in": emu_to_in(shape.left),
        "top_in": emu_to_in(shape.top),
        "width_in": emu_to_in(shape.width),
        "height_in": emu_to_in(shape.height),
    }
    if shape.is_placeholder:
        info["placeholder_type"] = str(shape.placeholder_format.type)

    if shape.has_text_frame and shape.text_frame.text.strip():
        info["text"] = shape.text_frame.text.strip()[:120]
        info["fonts"] = extract_font_info(shape.text_frame)

    if shape.has_table:
        t = shape.table
        info["table"] = {"rows": len(t.rows), "cols": len(t.columns)}

    if shape.has_chart:
        info["chart_type"] = str(shape.chart.chart_type)

    try:
        if shape.fill.type is not None:
            info["fill"] = color_to_hex(shape.fill.fore_color)
    except Exception:
        pass

    return info


def analyze(path: Path) -> dict:
    prs = Presentation(str(path))
    result = {
        "file": str(path),
        "slide_width_in": emu_to_in(prs.slide_width),
        "slide_height_in": emu_to_in(prs.slide_height),
        "aspect_ratio": round(prs.slide_width / prs.slide_height, 3),
        "theme_colors": extract_theme_colors(prs),
        "slide_count": len(prs.slides),
        "slides": [],
    }

    for idx, slide in enumerate(prs.slides, start=1):
        slide_info = {
            "index": idx,
            "layout_name": slide.slide_layout.name,
            "shape_count": len(slide.shapes),
            "shapes": [extract_shape(s) for s in slide.shapes],
        }
        try:
            bg = slide.background
            if bg.fill.type is not None:
                slide_info["background_fill"] = color_to_hex(bg.fill.fore_color)
        except Exception:
            pass
        result["slides"].append(slide_info)

    return result


def main() -> None:
    parser = argparse.ArgumentParser(description="레퍼런스 PPT 구조 분석")
    parser.add_argument("--input", required=True, help="레퍼런스 .pptx 경로")
    parser.add_argument("--output", help="출력 JSON 경로 (생략 시 표준출력)")
    args = parser.parse_args()

    result = analyze(Path(args.input))
    text = json.dumps(result, ensure_ascii=False, indent=2, default=str)

    if args.output:
        Path(args.output).write_text(text, encoding="utf-8")
        print(f"[web-ppt-generator] 분석 완료 -> {args.output}")
    else:
        print(text)


if __name__ == "__main__":
    main()
